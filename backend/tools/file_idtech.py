import os
import json
import zipfile
import chardet
import fitz  # PyMuPDF
from typing import List, Dict, Optional, Tuple
from pathlib import Path
import docx
import pptx
from bs4 import BeautifulSoup  # For HTML/XML
import pandas as pd  # For Excel/CSV
from azure.search.documents import SearchClient
from azure.search.documents.indexes import SearchIndexClient
from azure.search.documents.indexes.models import (
    SearchIndex, SimpleField, SearchFieldDataType, 
    SearchableField, VectorSearchProfile, 
    HnswAlgorithmConfiguration, VectorSearch,
    SearchField, SemanticConfiguration,
    SemanticPrioritizedFields, SemanticField
)
from azure.identity import DefaultAzureCredential
from openai import AzureOpenAI

class GameModdingKnowledgeBase:
    def __init__(self, azure_openai_endpoint: str, azure_search_endpoint: str, index_name: str):
        self.credential = DefaultAzureCredential()
        self.openai_client = AzureOpenAI(
            api_version="2023-05-15",
            azure_endpoint=azure_openai_endpoint,
            credential=self.credential
        )
        self.search_client = SearchClient(
            endpoint=azure_search_endpoint,
            index_name=index_name,
            credential=self.credential
        )
        self.index_client = SearchIndexClient(
            endpoint=azure_search_endpoint,
            credential=self.credential
        )
        self.index_name = index_name
        
        # Game modding specific configurations
        self.game_engines = ["mohaa", "idtech3", "gtkradiant"]
        self.file_types = {
            "code": [".c", ".cpp", ".h", ".py", ".js", ".java", ".sh", ".bat"],
            "scripts": [".scr", ".cfg", ".def", ".shader", ".script"],
            "docs": [".txt", ".md", ".doc", ".docx", ".pdf"],
            "assets": [".map", ".bsp", ".md3", ".pk3", ".wal", ".tga"],
            "config": [".json", ".xml", ".ini", ".cfg"],
            "subtitles": [".srt", ".str"]
        }

    def detect_game_engine(self, content: str) -> str:
        """Detect game engine from content using GPT"""
        try:
            response = self.openai_client.chat.completions.create(
                model="gpt-4",
                messages=[{
                    "role": "system",
                    "content": "Identify which game engine this content relates to (mohaa, idtech3, gtkradiant). Just return the engine name."
                }, {
                    "role": "user",
                    "content": content[:2000]  # First 2000 chars for detection
                }],
                temperature=0.1
            )
            engine = response.choices[0].message.content.lower()
            return engine if engine in self.game_engines else "unknown"
        except:
            return "unknown"

    def extract_metadata(self, file_path: str, content: str) -> Dict:
        """Extract metadata from file path and content"""
        file_ext = Path(file_path).suffix.lower()
        file_type = next(
            (k for k, v in self.file_types.items() if file_ext in v),
            "other"
        )
        
        engine = self.detect_game_engine(content)
        
        return {
            "file_type": file_type,
            "engine": engine,
            "file_name": Path(file_path).name,
            "file_path": str(Path(file_path).relative_to(Path(file_path).anchor))
        }

    def read_file(self, file_path: str) -> Tuple[str, Dict]:
        """Read any supported file type and return content + metadata"""
        ext = Path(file_path).suffix.lower()
        
        # Text-based files
        if ext in {".txt", ".md", ".cfg", ".ini", ".shader", ".script", ".scr", ".srt", ".str"}:
            with open(file_path, 'rb') as f:
                raw = f.read()
                encoding = chardet.detect(raw)['encoding'] or 'utf-8'
                content = raw.decode(encoding, errors='ignore')
        
        # PDF files
        elif ext == ".pdf":
            content = []
            with fitz.open(file_path) as doc:
                for page in doc:
                    content.append(page.get_text())
            content = "\n".join(content)
        
        # Word documents
        elif ext == ".docx":
            doc = docx.Document(file_path)
            content = "\n".join([para.text for para in doc.paragraphs])
        
        # PowerPoint
        elif ext == ".pptx":
            prs = pptx.Presentation(file_path)
            content = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        content.append(shape.text)
            content = "\n".join(content)
        
        # HTML/XML
        elif ext in {".html", ".htm", ".xml"}:
            with open(file_path, 'r', encoding='utf-8') as f:
                soup = BeautifulSoup(f.read(), 'html.parser')
                content = soup.get_text()
        
        # JSON
        elif ext == ".json":
            with open(file_path, 'r', encoding='utf-8') as f:
                data = json.load(f)
                content = json.dumps(data, indent=2)
        
        # CSV/Excel
        elif ext in {".csv", ".xlsx", ".xls"}:
            df = pd.read_excel(file_path) if ext in {".xlsx", ".xls"} else pd.read_csv(file_path)
            content = df.to_string()
        
        # PK3/ZIP files (common in idTech3)
        elif ext == ".pk3" or ext == ".zip":
            content = []
            with zipfile.ZipFile(file_path) as z:
                for name in z.namelist():
                    if not name.endswith('/'):  # Skip directories
                        try:
                            with z.open(name) as f:
                                content.append(f"{name}:\n{f.read().decode('utf-8', errors='ignore')}")
                        except:
                            continue
            content = "\n".join(content)
        
        else:
            raise ValueError(f"Unsupported file type: {ext}")
        
        metadata = self.extract_metadata(file_path, content)
        return content, metadata

    def chunk_content(self, content: str, metadata: Dict, chunk_size: int = 1500, overlap: int = 300) -> List[Dict]:
        """Chunk content with context-aware splitting"""
        # Special handling for code/scripts
        if metadata["file_type"] in ["code", "scripts"]:
            return self._chunk_code(content, metadata, chunk_size)
        
        # Standard text chunking with overlap
        chunks = []
        start = 0
        length = len(content)
        
        while start < length:
            end = start + chunk_size
            chunk = content[start:end]
            
            # Try to end at a paragraph break if possible
            last_newline = chunk.rfind('\n')
            if last_newline > 0 and (end - last_newline) < 200:  # Don't go too far back
                end = start + last_newline
                chunk = content[start:end]
            
            chunks.append({
                "content": chunk,
                "chunk_num": len(chunks) + 1,
                "start_pos": start,
                "end_pos": end,
                **metadata
            })
            
            start = end - overlap if (end - overlap) > start else end
        
        return chunks

    def _chunk_code(self, content: str, metadata: Dict, chunk_size: int) -> List[Dict]:
        """Specialized chunking for code files"""
        lines = content.split('\n')
        chunks = []
        current_chunk = []
        current_length = 0
        
        for i, line in enumerate(lines):
            line_length = len(line) + 1  # +1 for newline
            
            if current_length + line_length > chunk_size and current_chunk:
                # Add current chunk before it gets too big
                chunks.append({
                    "content": '\n'.join(current_chunk),
                    "chunk_num": len(chunks) + 1,
                    "start_line": i - len(current_chunk) + 1,
                    "end_line": i,
                    **metadata
                })
                current_chunk = []
                current_length = 0
            
            current_chunk.append(line)
            current_length += line_length
        
        # Add the last chunk if there's anything left
        if current_chunk:
            chunks.append({
                "content": '\n'.join(current_chunk),
                "chunk_num": len(chunks) + 1,
                "start_line": len(lines) - len(current_chunk) + 1,
                "end_line": len(lines),
                **metadata
            })
        
        return chunks

    def process_folder(self, folder_path: str) -> List[Dict]:
        """Process all supported files in a folder"""
        documents = []
        
        for root, _, files in os.walk(folder_path):
            for file in files:
                file_path = os.path.join(root, file)
                try:
                    content, metadata = self.read_file(file_path)
                    chunks = self.chunk_content(content, metadata)
                    documents.extend(chunks)
                    print(f"Processed {file_path} ({len(chunks)} chunks)")
                except ValueError as e:
                    print(f"Skipping {file_path}: {str(e)}")
                except Exception as e:
                    print(f"Error processing {file_path}: {str(e)}")
        
        return documents

    def create_index(self):
        """Create Azure AI Search index optimized for game modding content"""
        fields = [
            SimpleField(name="id", type=SearchFieldDataType.String, key=True),
            SearchableField(name="content", type=SearchFieldDataType.String),
            SearchableField(name="file_name", type=SearchFieldDataType.String, filterable=True),
            SearchableField(name="file_path", type=SearchFieldDataType.String),
            SearchableField(name="file_type", type=SearchFieldDataType.String, filterable=True, facetable=True),
            SearchableField(name="engine", type=SearchFieldDataType.String, filterable=True, facetable=True),
            SimpleField(name="chunk_num", type=SearchFieldDataType.Int32),
            SimpleField(name="start_pos", type=SearchFieldDataType.Int32),
            SimpleField(name="end_pos", type=SearchFieldDataType.Int32),
            SimpleField(name="start_line", type=SearchFieldDataType.Int32),
            SimpleField(name="end_line", type=SearchFieldDataType.Int32),
            SearchField(
                name="content_vector",
                type=SearchFieldDataType.Collection(SearchFieldDataType.Single),
                searchable=True,
                vector_search_dimensions=1536,
                vector_search_profile_name="my-vector-profile"
            )
        ]

        vector_search = VectorSearch(
            profiles=[VectorSearchProfile(
                name="my-vector-profile",
                algorithm_configuration_name="my-algorithms-config",
                vectorizer=None
            )],
            algorithms=[HnswAlgorithmConfiguration(
                name="my-algorithms-config",
                parameters={
                    "m": 4,
                    "efConstruction": 400,
                    "efSearch": 500,
                    "metric": "cosine"
                }
            )]
        )

        semantic_config = SemanticConfiguration(
            name="my-semantic-config",
            prioritized_fields=SemanticPrioritizedFields(
                title_field=SemanticField(field_name="file_name"),
                content_fields=[SemanticField(field_name="content")],
                keywords_fields=[SemanticField(field_name="file_type")]
            )
        )

        index = SearchIndex(
            name=self.index_name,
            fields=fields,
            vector_search=vector_search,
            semantic_settings={"configurations": [semantic_config]}
        )

        try:
            self.index_client.create_index(index)
            print(f"Index {self.index_name} created successfully")
        except Exception as e:
            print(f"Error creating index: {e}")

    def upload_documents(self, documents: List[Dict]):
        """Upload documents with embeddings"""
        # Generate embeddings in batches
        batch_size = 10
        for i in range(0, len(documents), batch_size):
            batch = documents[i:i+batch_size]
            
            # Generate embeddings
            embeddings = self.openai_client.embeddings.create(
                input=[doc["content"] for doc in batch],
                model="text-embedding-ada-002"
            )
            
            # Add embeddings to documents
            for j, doc in enumerate(batch):
                doc["id"] = f"{doc['file_path']}_{doc['chunk_num']}"
                doc["content_vector"] = embeddings.data[j].embedding
            
            # Upload batch
            try:
                result = self.search_client.upload_documents(documents=batch)
                print(f"Uploaded batch {i//batch_size + 1}")
            except Exception as e:
                print(f"Error uploading batch: {e}")

    def search(self, query: str, filters: Optional[str] = None, n: int = 5) -> List[Dict]:
        """Search with optional filters (e.g., "engine eq 'mohaa'")"""
        # Generate query embedding
        embedding = self.openai_client.embeddings.create(
            input=query,
            model="text-embedding-ada-002"
        ).data[0].embedding

        # Build search parameters
        search_params = {
            "search_text": query,
            "vector_queries": [{
                "text": query,
                "fields": "content_vector",
                "k_nearest_neighbors": n * 10,  # Oversampling
                "exhaustive": True
            }],
            "query_type": "semantic",
            "semantic_configuration_name": "my-semantic-config",
            "select": "content,file_name,file_path,file_type,engine",
            "top": n,
            "filter": filters
        }

        results = []
        for r in self.search_client.search(**search_params):
            results.append({
                "content": r["content"],
                "file_name": r["file_name"],
                "file_path": r["file_path"],
                "file_type": r["file_type"],
                "engine": r["engine"],
                "score": r.get("@search.score", 0),
                "reranker_score": r.get("@search.reranker_score", 0)
            })
        
        return results